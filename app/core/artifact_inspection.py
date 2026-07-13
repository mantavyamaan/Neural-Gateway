"""
Artifact Inspection Layer.

Reads real file metadata (PyMuPDF / Pillow / ffprobe / openpyxl /
python-pptx) when a physical file is given, or falls back to
prompt-keyword heuristics for backward-compatible / synthetic-test usage.
Also does lightweight semantic extraction (short text preview + topic
inference) and prompt/upload conflict detection.
"""

import json
import mimetypes
import os
import re
import shutil
import subprocess
from typing import Any, Dict, List, Optional

from app.models.schemas import ArtifactProfile

# ---- Optional dependency imports (graceful degradation) -----------------


def _try_import(name: str):
    try:
        return __import__(name)
    except Exception:
        return None


fitz = _try_import("fitz")          # PyMuPDF -> PDFs
openpyxl = _try_import("openpyxl")  # Excel
pptx_mod = _try_import("pptx")      # python-pptx -> PowerPoint
mutagen = _try_import("mutagen")    # audio metadata

from typing import Any, Dict, List, Optional
try:
    from PIL import Image as _PILImage
except Exception:
    _PILImage = None  # type: ignore[assignment]


def _ffprobe_available() -> bool:
    return shutil.which("ffprobe") is not None


ARTIFACT_LIBS = {
    "pymupdf": fitz is not None,
    "pillow": _PILImage is not None,
    "openpyxl": openpyxl is not None,
    "python-pptx": pptx_mod is not None,
    "mutagen": mutagen is not None,
    "ffprobe": _ffprobe_available(),
}


# ---- Format detection -----------------------------------------------------

EXTENSION_FORMAT_MAP: Dict[str, str] = {
    ".pdf": "pdf",
    ".png": "image", ".jpg": "image", ".jpeg": "image",
    ".gif": "image", ".bmp": "image", ".tiff": "image", ".webp": "image",
    ".mp3": "audio", ".wav": "audio", ".m4a": "audio",
    ".flac": "audio", ".ogg": "audio", ".aac": "audio",
    ".mp4": "video", ".mov": "video", ".mkv": "video",
    ".avi": "video", ".webm": "video",
    ".xlsx": "spreadsheet", ".xls": "spreadsheet", ".csv": "spreadsheet",
    ".pptx": "presentation", ".ppt": "presentation",
    ".txt": "text", ".md": "text",
}


def detect_format(path: str) -> Optional[str]:
    """Determine the canonical SUPPORTED_FORMATS label for a file path.

    Falls back to MIME-type sniffing when the extension is unknown.
    Returns None if the format cannot be mapped (caller decides how to handle).
    """
    ext = os.path.splitext(path)[1].lower()
    if ext in EXTENSION_FORMAT_MAP:
        return EXTENSION_FORMAT_MAP[ext]
    mime, _ = mimetypes.guess_type(path)
    if mime:
        if mime.startswith("image/"):
            return "image"
        if mime.startswith("audio/"):
            return "audio"
        if mime.startswith("video/"):
            return "video"
        if mime == "application/pdf":
            return "pdf"
        if "spreadsheet" in mime or mime == "text/csv":
            return "spreadsheet"
        if "presentation" in mime:
            return "presentation"
        if mime.startswith("text/"):
            return "text"
    return None


# ---- Per-format inspectors --------------------------------------------------

def _inspect_pdf(path: str, profile: ArtifactProfile) -> None:
    """Read PDF structure with PyMuPDF; derive scan/text/table signals."""
    if fitz is None:
        return
    try:
        doc = fitz.open(path)
        profile.page_count = doc.page_count
        first = doc.load_page(0)
        text = first.get_text("text") or ""
        images = first.get_images(full=True)
        char_count = len(text.strip())
        profile.text_density = min(1.0, char_count / 1500.0)
        profile.scan_likelihood = 0.85 if (char_count < 100 and images) else 0.10
        profile.handwriting_likelihood = 0.05
        profile.table_density = 0.70 if text.count("\t") > 10 else 0.10
        profile.detected_language = "en"  # replace with langdetect in prod
        doc.close()
    except Exception:
        pass


def _inspect_image(path: str, profile: ArtifactProfile) -> None:
    if _PILImage is None:
        return
    try:
        img = _PILImage.open(path)
        w, h = img.size
        profile.scan_likelihood = 0.60 if h > w else 0.20
        profile.handwriting_likelihood = 0.10
        profile.table_density = 0.10
        profile.chart_density = 0.10
        profile.detected_language = "en"
    except Exception:
        pass


def _inspect_media_ffprobe(path: str, profile: ArtifactProfile, is_video: bool) -> bool:
    """Use ffprobe to read duration and stream info. Returns True on success."""
    if not _ffprobe_available():
        return False
    try:
        out = subprocess.run(
            ["ffprobe", "-v", "quiet", "-print_format", "json",
             "-show_format", "-show_streams", path],
            capture_output=True, text=True, timeout=20, shell=False
        )
        meta = json.loads(out.stdout or "{}")
        duration = float(meta.get("format", {}).get("duration", 0)) or None
        has_audio = any(s.get("codec_type") == "audio" for s in meta.get("streams", []))
        if is_video:
            profile.video_duration_sec = int(duration) if duration else None
            profile.audio_quality = 0.70 if has_audio else 0.0
            profile.chart_density = 0.20
        else:
            profile.audio_duration_sec = int(duration) if duration else None
            profile.audio_quality = 0.75
        profile.detected_language = "en"
        return True
    except Exception:
        return False


def _inspect_audio(path: str, profile: ArtifactProfile) -> None:
    if _inspect_media_ffprobe(path, profile, is_video=False):
        return
    if mutagen is not None:
        try:
            audio = mutagen.File(path)
            if audio is not None and audio.info is not None:
                profile.audio_duration_sec = int(audio.info.length)
                profile.audio_quality = 0.75
                profile.detected_language = "en"
        except Exception:
            pass


def _inspect_video(path: str, profile: ArtifactProfile) -> None:
    _inspect_media_ffprobe(path, profile, is_video=True)


def _inspect_spreadsheet(path: str, profile: ArtifactProfile) -> None:
    if openpyxl is None:
        return
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        n_sheets = len(wb.sheetnames)
        max_cells = 0
        for ws in wb.worksheets:
            max_cells = max(max_cells, (ws.max_row or 0) * (ws.max_column or 0))
        profile.spreadsheet_complexity = min(1.0, 0.3 + 0.2 * n_sheets + max_cells / 50000.0)
        profile.table_density = 0.90
        profile.detected_language = "en"
        wb.close()
    except Exception:
        pass


def _inspect_presentation(path: str, profile: ArtifactProfile) -> None:
    if pptx_mod is None:
        return
    try:
        from pptx import Presentation
        prs = Presentation(path)
        n_slides = len(prs.slides)
        profile.presentation_complexity = min(1.0, 0.3 + n_slides / 40.0)
        profile.chart_density = 0.55
        profile.detected_language = "en"
    except Exception:
        pass


# ---- Lightweight semantic extraction (text preview + topic) ---------------

TOPIC_KEYWORDS: Dict[str, List[str]] = {
    "legal_contract": ["agreement", "obligations", "liability", "indemnif",
                       "nda", "confidential", "party of the", "hereinafter", "clause"],
    "financial_document": ["invoice", "balance sheet", "revenue", "ebitda",
                           "statement", "tax", "amount due", "fiscal"],
    "medical_record": ["patient", "diagnosis", "treatment", "prescription",
                       "symptom", "clinical"],
    "research_paper": ["abstract", "we propose", "related work", "experiment",
                       "benchmark", "et al", "references"],
    "security_report": ["vulnerability", "cve", "exploit", "payload", "attack surface"],
    "support_record": ["ticket", "customer", "refund", "complaint", "resolution"],
}


def extract_text_preview(path: str, fmt: str, max_words: int = 500) -> Optional[str]:
    """Pull a small text preview from a file for topic inference. Best-effort."""
    try:
        if fmt == "pdf" and fitz is not None:
            doc = fitz.open(path)
            text = doc.load_page(0).get_text("text") or ""
            doc.close()
        elif fmt == "spreadsheet" and openpyxl is not None:
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            ws = wb.active
            cells: List[Any] = []
            for row in ws.iter_rows(max_row=20, values_only=True):
                cells.extend(str(c) for c in row if c is not None)
            text = " ".join(cells)
            wb.close()
        elif fmt == "presentation" and pptx_mod is not None:
            from pptx import Presentation
            prs = Presentation(path)
            chunks = []
            for slide in list(prs.slides)[:3]:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        chunks.append(shape.text_frame.text)
            text = " ".join(chunks)
        elif fmt == "text":
            with open(path, "r", errors="ignore") as fh:
                text = fh.read(8000)
        else:
            return None
        return " ".join(text.split()[:max_words]) or None
    except Exception:
        return None


def infer_topic(preview: Optional[str]) -> Optional[str]:
    """Score the preview against TOPIC_KEYWORDS; return best match or None."""
    if not preview:
        return None
    low = preview.lower()
    best_topic, best_score = None, 0
    for topic, kws in TOPIC_KEYWORDS.items():
        score = sum(1 for kw in kws if re.search(r'(?<!\w)' + re.escape(kw) + r'(?!\w)', low, re.IGNORECASE))
        if score > best_score:
            best_topic, best_score = topic, score
    return best_topic if best_score >= 2 else None


# ---- Prompt <-> upload conflict detection ---------------------------------

PROMPT_MODALITY_HINTS: Dict[str, List[str]] = {
    "image": ["this image", "this picture", "this photo", "the photo"],
    "pdf": ["this pdf", "this document", "the document", "this contract"],
    "audio": ["this recording", "this audio", "this call", "the meeting recording"],
    "video": ["this video", "this clip", "the footage"],
    "spreadsheet": ["this spreadsheet", "this excel", "this workbook"],
    "presentation": ["this slide deck", "this presentation", "these slides"],
}


def detect_conflicts(prompt: str, profiles: List[ArtifactProfile]) -> List[str]:
    """Flag mismatches between the modality the prompt implies and what was uploaded.

    Resolution policy: trust the uploaded artifact, log a warning flag. The
    router keeps routing on real artifacts; the flag surfaces in the audit
    record so a reviewer (or upstream UX) can prompt the user if desired.
    """
    flags: List[str] = []
    p = prompt.lower()
    present_formats = {prof.format for prof in profiles}
    for modality, phrases in PROMPT_MODALITY_HINTS.items():
        if any(phrase in p for phrase in phrases) and modality not in present_formats:
            uploaded = sorted(present_formats - {"text"}) or ["none"]
            flags.append(
                f"prompt_implies_{modality}_but_uploaded_{'/'.join(uploaded)}"
            )
    return flags


# ---- Main entry point -------------------------------------------------------

def inspect_artifacts(
    input_formats: Optional[List[str]] = None,
    prompt: str = "",
    artifact_hints: Optional[List[Dict[str, Any]]] = None,
    files: Optional[List[str]] = None,
) -> List[ArtifactProfile]:
    """Produce one ArtifactProfile per input.

    Two modes:
      1. FILE MODE (preferred): pass `files=[...]`. Real metadata is read from
         each file using PyMuPDF/Pillow/FFprobe/openpyxl/python-pptx when those
         libraries are available, with graceful degradation otherwise.
      2. LEGACY MODE: pass `input_formats=[...]`. Falls back to prompt-keyword
         heuristics. Preserved for backward compatibility and synthetic tests.

    Caller-supplied `artifact_hints` always win over both file reads and
    heuristics (ground truth injection from an upstream system).
    """
    artifact_hints = artifact_hints or []
    profiles: List[ArtifactProfile] = []

    # ---------------- FILE MODE ----------------
    if files:
        hints_by_format = {h.get("format"): h for h in artifact_hints if "format" in h}
        for path in files:
            fmt = detect_format(path)
            if fmt is None:
                profiles.append(ArtifactProfile(format="text"))
                continue
            profile = ArtifactProfile(format=fmt)
            try:
                if fmt == "pdf":
                    _inspect_pdf(path, profile)
                elif fmt == "image":
                    _inspect_image(path, profile)
                elif fmt == "audio":
                    _inspect_audio(path, profile)
                elif fmt == "video":
                    _inspect_video(path, profile)
                elif fmt == "spreadsheet":
                    _inspect_spreadsheet(path, profile)
                elif fmt == "presentation":
                    _inspect_presentation(path, profile)
            except Exception:
                pass
            try:
                profile.source_path = path
                profile.file_size_bytes = os.path.getsize(path)
                profile.extraction_method = "library"
                preview = extract_text_preview(path, fmt)
                profile.extracted_text_preview = preview
                profile.inferred_topic = infer_topic(preview)
            except Exception:
                pass
            h = hints_by_format.get(fmt, {})
            for k, v in h.items():
                if k != "format" and hasattr(profile, k):
                    setattr(profile, k, v)
            profiles.append(profile)
        return profiles

    # ---------------- LEGACY HEURISTIC MODE ----------------
    input_formats = input_formats or ["text"]
    p = prompt.lower()
    hints_by_format = {h.get("format"): h for h in artifact_hints if "format" in h}
    for fmt in sorted(set(input_formats)):
        h = hints_by_format.get(fmt, {})
        profile = ArtifactProfile(format=fmt)
        if fmt == "pdf":
            profile.page_count = h.get("page_count", 30 if "contract" in p else 10)
            profile.text_density = h.get("text_density",
                0.25 if any(k in p for k in ["scanned", "scan", "invoice"]) else 0.80)
            profile.scan_likelihood = h.get("scan_likelihood",
                0.85 if any(k in p for k in ["scanned", "scan", "invoice"]) else 0.10)
            profile.handwriting_likelihood = h.get("handwriting_likelihood",
                0.50 if "handwritten" in p else 0.05)
            profile.table_density = h.get("table_density",
                0.70 if any(k in p for k in ["invoice", "statement", "table"]) else 0.10)
            profile.chart_density = h.get("chart_density",
                0.55 if any(k in p for k in ["chart", "graph"]) else 0.10)
            profile.detected_language = h.get("detected_language", "en")
        elif fmt == "image":
            profile.scan_likelihood = h.get("scan_likelihood",
                0.80 if any(k in p for k in ["extract text", "scan", "ocr"]) else 0.10)
            profile.handwriting_likelihood = h.get("handwriting_likelihood",
                0.60 if "handwritten" in p else 0.10)
            profile.table_density = h.get("table_density",
                0.65 if any(k in p for k in ["invoice", "receipt", "table"]) else 0.10)
            profile.chart_density = h.get("chart_density",
                0.50 if "chart" in p else 0.05)
            profile.detected_language = h.get("detected_language", "en")
        elif fmt == "audio":
            profile.audio_duration_sec = h.get("audio_duration_sec",
                900 if any(k in p for k in ["call", "meeting"]) else 120)
            profile.audio_quality = h.get("audio_quality", 0.75)
            profile.detected_language = h.get("detected_language", "en")
        elif fmt == "video":
            profile.video_duration_sec = h.get("video_duration_sec", 600)
            profile.audio_quality = h.get("audio_quality", 0.70)
            profile.chart_density = h.get("chart_density", 0.20)
            profile.detected_language = h.get("detected_language", "en")
        elif fmt == "spreadsheet":
            profile.spreadsheet_complexity = h.get("spreadsheet_complexity",
                0.75 if any(k in p for k in ["model", "forecast", "multi-she"]) else 0.40)
            profile.table_density = h.get("table_density", 0.90)
            profile.detected_language = h.get("detected_language", "en")
        elif fmt == "presentation":
            profile.presentation_complexity = h.get("presentation_complexity", 0.50)
            profile.chart_density = h.get("chart_density", 0.55)
            profile.detected_language = h.get("detected_language", "en")
        else:
            profile.detected_language = h.get("detected_language", "en")
        profiles.append(profile)
    return profiles
